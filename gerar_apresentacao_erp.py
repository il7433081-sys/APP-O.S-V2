#!/usr/bin/env python3
"""Gera apresentação PowerPoint — ERP da Oficina (texto de apresentação)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Identidade visual (mesma paleta do App O.S. Digital)
NAVY = RGBColor(0, 56, 168)
NAVY_DARK = RGBColor(0, 40, 120)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(80, 80, 80)
LIGHT_BG = RGBColor(232, 238, 248)
ACCENT = RGBColor(0, 120, 200)

SAIDA = Path(__file__).resolve().parent.parent.parent / "Apresentacao_ERP_Oficina.pptx"

SLIDES: list[dict] = [
    {
        "tipo": "titulo",
        "titulo": "Sistema ERP da Oficina",
        "subtitulo": "Planejamento de Recursos Empresariais\nIntegração Escritório + Pátio via API",
        "notas": (
            "Bom dia a todos. Hoje eu tenho o prazer de apresentar a vocês o sistema "
            "que desenvolvi e que vamos implementar na oficina."
        ),
    },
    {
        "tipo": "conteudo",
        "titulo": "Uma solução feita para nós",
        "bullets": [
            "Sistema desenvolvido internamente para a rotina da oficina náutica.",
            "Implementação pensada para escritório e pátio trabalharem juntos.",
            "Tecnologia de ponta aplicada ao nosso dia a dia operacional.",
        ],
        "notas": (
            "Bom dia a todos. Hoje eu tenho o prazer de apresentar a vocês o sistema "
            "que desenvolvi e que vamos implementar na oficina."
        ),
    },
    {
        "tipo": "destaque",
        "titulo": "O que é um ERP?",
        "destaque": "Enterprise Resource Planning",
        "subdestaque": "Planejamento de Recursos Empresariais",
        "texto": (
            "Na engenharia de software, isso não é apenas um aplicativo comum, "
            "mas sim um sistema integrado completo feito para gerir todas as "
            "operações da empresa em um só lugar."
        ),
        "notas": (
            "Para que todos entendam o nível da solução que estamos trazendo para "
            "o negócio, o que temos aqui é um legítimo ERP, que vem do inglês "
            "Enterprise Resource Planning e significa Planejamento de Recursos "
            "Empresariais. Na engenharia de software, isso não é apenas um "
            "aplicativo comum, mas sim um sistema integrado completo feito para "
            "gerir todas as operações da empresa em um só lugar."
        ),
    },
    {
        "tipo": "conteudo",
        "titulo": "Tudo integrado em um só lugar",
        "bullets": [
            "Histórico completo das embarcações e motores.",
            "Regras operacionais da oficina (comissões, tabelas, revisões).",
            "Ordens de serviço, estoque, pré-orçamentos e agendamentos.",
            "Cadastros, permissões e relatórios centralizados.",
        ],
        "notas": (
            "...desde o histórico das embarcações até as nossas regras operacionais, "
            "como comissões e tabelas de revisões."
        ),
    },
    {
        "tipo": "secao",
        "titulo": "O grande diferencial técnico",
        "subtitulo": "Integração via API",
        "notas": (
            "Mas o grande diferencial técnico deste projeto é como ele trabalha. "
            "Para fazer o computador do escritório e o tablet do pátio conversarem, "
            "eu integrei este ERP com uma API."
        ),
    },
    {
        "tipo": "destaque",
        "titulo": "O que é uma API?",
        "destaque": "Application Programming Interface",
        "subdestaque": "Interface de Programação de Aplicações",
        "texto": (
            "Uma ponte invisível de bastidores que conecta os sistemas "
            "de forma segura, rápida e organizada."
        ),
        "notas": (
            "API significa Application Programming Interface, ou Interface de "
            "Programação de Aplicações."
        ),
    },
    {
        "tipo": "conteudo",
        "titulo": "API na prática — o que ela faz",
        "bullets": [
            "Dado inserido no pátio → viaja instantaneamente para a base central.",
            "Comunicação segura entre escritório e tablet do mecânico.",
            "Códigos separados, dados sincronizados — tudo em perfeita harmonia.",
            "Segurança e velocidade para a rotina da oficina.",
        ],
        "notas": (
            "Na prática, a API funciona como uma ponte invisível de bastidores. "
            "Ela garante que quando um dado for inserido lá no pátio, ele viaje "
            "de forma segura e instantânea para a nossa base central, sem misturar "
            "os códigos e mantendo tudo em perfeita harmonia. É o que há de melhor "
            "em engenharia de software para garantir segurança e velocidade para "
            "a nossa rotina."
        ),
    },
    {
        "tipo": "diagrama",
        "titulo": "Fluxo de integração",
        "notas": "Resumo visual do fluxo escritório ↔ API ↔ pátio.",
    },
    {
        "tipo": "secao",
        "titulo": "Hora de ver na prática",
        "subtitulo": "Demonstração do sistema",
        "notas": (
            "Sabendo disso, eu quero deixar a teoria de lado e mostrar para vocês, "
            "na prática, como essa tecnologia vai rodar na nossa tela. "
            "Vamos dar uma olhada no programa..."
        ),
    },
    {
        "tipo": "fechamento",
        "titulo": "Obrigado!",
        "subtitulo": "Perguntas?",
        "notas": "Encerramento da apresentação — abrir para dúvidas.",
    },
]


def _caixa_fundo(slide, cor: RGBColor, topo: float = 0, altura: float | None = None) -> None:
    alt = altura or Inches(7.5)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(topo), Inches(10), alt)  # retângulo
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def _texto(slide, left, top, width, height, texto, *, tamanho=18, negrito=False, cor=GRAY, alinhamento=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = negrito
    p.font.color.rgb = cor
    p.alignment = alinhamento
    return box


def _slide_titulo(prs: Presentation, dados: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _caixa_fundo(slide, NAVY, altura=Inches(7.5))

    _texto(
        slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.2),
        dados["titulo"], tamanho=40, negrito=True, cor=WHITE, alinhamento=PP_ALIGN.CENTER,
    )
    _texto(
        slide, Inches(1.2), Inches(3.5), Inches(7.6), Inches(1.5),
        dados.get("subtitulo", ""), tamanho=20, cor=RGBColor(200, 220, 255), alinhamento=PP_ALIGN.CENTER,
    )

    barra = slide.shapes.add_shape(1, Inches(3.5), Inches(5.2), Inches(3), Inches(0.06))
    barra.fill.solid()
    barra.fill.fore_color.rgb = WHITE
    barra.line.fill.background()

    if dados.get("notas"):
        slide.notes_slide.notes_text_frame.text = dados["notas"]


def _slide_conteudo(prs: Presentation, dados: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    faixa = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.35))
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = NAVY
    faixa.line.fill.background()

    _texto(
        slide, Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.8),
        dados["titulo"], tamanho=28, negrito=True, cor=WHITE,
    )

    y = Inches(1.8)
    for item in dados.get("bullets", []):
        box = slide.shapes.add_textbox(Inches(0.9), y, Inches(8.2), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
        p.level = 0
        p.space_before = Pt(4)
        # marcador manual via bullet unicode
        p.text = "▸  " + item
        y += Inches(0.85)

    if dados.get("notas"):
        slide.notes_slide.notes_text_frame.text = dados["notas"]


def _slide_destaque(prs: Presentation, dados: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _caixa_fundo(slide, LIGHT_BG)

    _texto(slide, Inches(0.8), Inches(0.8), Inches(8.4), Inches(0.7), dados["titulo"], tamanho=26, negrito=True, cor=NAVY, alinhamento=PP_ALIGN.CENTER)

    card = slide.shapes.add_shape(1, Inches(1.2), Inches(2.0), Inches(7.6), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY
    card.line.fill.background()

    _texto(slide, Inches(1.5), Inches(2.35), Inches(7.0), Inches(0.8), dados.get("destaque", ""), tamanho=24, negrito=True, cor=WHITE, alinhamento=PP_ALIGN.CENTER)
    _texto(slide, Inches(1.5), Inches(3.15), Inches(7.0), Inches(0.7), dados.get("subdestaque", ""), tamanho=18, cor=RGBColor(180, 210, 255), alinhamento=PP_ALIGN.CENTER)

    _texto(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(1.5), dados.get("texto", ""), tamanho=18, cor=GRAY, alinhamento=PP_ALIGN.CENTER)

    if dados.get("notas"):
        slide.notes_slide.notes_text_frame.text = dados["notas"]


def _slide_secao(prs: Presentation, dados: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _caixa_fundo(slide, NAVY_DARK)

    _texto(slide, Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.0), dados["titulo"], tamanho=36, negrito=True, cor=WHITE, alinhamento=PP_ALIGN.CENTER)
    _texto(slide, Inches(1.0), Inches(3.7), Inches(8.0), Inches(0.8), dados.get("subtitulo", ""), tamanho=22, cor=ACCENT, alinhamento=PP_ALIGN.CENTER)

    if dados.get("notas"):
        slide.notes_slide.notes_text_frame.text = dados["notas"]


def _slide_diagrama(prs: Presentation, dados: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    faixa = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = NAVY
    faixa.line.fill.background()
    _texto(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7), dados["titulo"], tamanho=26, negrito=True, cor=WHITE)

    blocos = [
        ("Escritório\n(computador)", NAVY),
        ("API\n(ponte segura)", ACCENT),
        ("Pátio\n(tablet)", NAVY),
    ]
    xs = [Inches(0.9), Inches(3.8), Inches(6.7)]
    for (label, cor), x in zip(blocos, xs):
        sh = slide.shapes.add_shape(1, x, Inches(2.8), Inches(2.2), Inches(1.6))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        _texto(slide, x + Inches(0.15), Inches(3.0), Inches(1.9), Inches(1.2), label, tamanho=16, negrito=True, cor=WHITE, alinhamento=PP_ALIGN.CENTER)

    for x_ini in (Inches(3.1), Inches(6.0)):
        seta = slide.shapes.add_shape(1, x_ini, Inches(3.45), Inches(0.7), Inches(0.25))
        seta.fill.solid()
        seta.fill.fore_color.rgb = GRAY
        seta.line.fill.background()

    _texto(
        slide, Inches(0.8), Inches(5.0), Inches(8.4), Inches(1.2),
        "Dado inserido no pátio  →  sincronizado em tempo real  →  visível no escritório",
        tamanho=17, cor=GRAY, alinhamento=PP_ALIGN.CENTER,
    )

    if dados.get("notas"):
        slide.notes_slide.notes_text_frame.text = dados["notas"]


def _slide_fechamento(prs: Presentation, dados: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _caixa_fundo(slide, NAVY)

    _texto(slide, Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.0), dados["titulo"], tamanho=44, negrito=True, cor=WHITE, alinhamento=PP_ALIGN.CENTER)
    _texto(slide, Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.8), dados.get("subtitulo", ""), tamanho=24, cor=RGBColor(200, 220, 255), alinhamento=PP_ALIGN.CENTER)

    if dados.get("notas"):
        slide.notes_slide.notes_text_frame.text = dados["notas"]


def gerar() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    handlers = {
        "titulo": _slide_titulo,
        "conteudo": _slide_conteudo,
        "destaque": _slide_destaque,
        "secao": _slide_secao,
        "diagrama": _slide_diagrama,
        "fechamento": _slide_fechamento,
    }

    for dados in SLIDES:
        tipo = dados.get("tipo", "conteudo")
        handlers[tipo](prs, dados)

    prs.save(SAIDA)
    return SAIDA


if __name__ == "__main__":
    caminho = gerar()
    print(f"Apresentação gerada: {caminho}")
